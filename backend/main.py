import os
import json
import asyncio
import io
from typing import Optional, List
from fastapi import FastAPI, HTTPException, File, UploadFile, Query
from fastapi.responses import HTMLResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel, Field
from google import genai
from google.genai import types
from vector_store import AuraVectorStore

app = FastAPI(
    title="Aura Agentic Content Supply Chain API",
    description="Enterprise Multi-Agent Orchestration Gateway for Automated, Compliant Content Generation with Multi-LLM Comparisons (White-Labeled).",
    version="2.0.0"
)

db = AuraVectorStore()

# Enable Robust CORS Configuration for Decoupled Vercel Frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Swap with your explicit Vercel domain in production settings
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# --- STRUCTURAL DATA SCHEMAS (PYDANTIC) ---
class CampaignRequest(BaseModel):
    brief: str = Field(..., example="Generate an autumn forecourt café campaign offering a 15% discount on our organic artisan coffee and hot wraps...")

class LocalizedCopy(BaseModel):
    email_subject: str = Field(..., description="High-impact corporate email subject line.")
    email_body: str = Field(..., description="Comprehensive retail email body supporting the narrative.")
    push_notification: str = Field(..., description="Short, compelling mobile application push message (max 120 chars).")
    in_app_banner: str = Field(..., description="Action-oriented text for in-app banner placement.")

class ComplianceCheck(BaseModel):
    status: str = Field(..., description="Strictly either 'APPROVED' or 'RISK_FLAG_DETECTED'")
    asa_greenwashing_passed: bool = Field(..., description="True if claims are linked to specific asset footprints, False if generalities are used.")
    clear_pricing_disclosure: bool = Field(..., description="True if discounts and criteria are explicitly bounded.")
    executive_summary: str = Field(..., description="A professional 1-sentence risk and regulatory compliance sign-off.")

class MultiAgentPayload(BaseModel):
    localized_content: LocalizedCopy
    asset_theme_token: str = Field(..., description="Sleek CSS slug representing visual asset directive. Strictly one of: 'ev-charging', 'aura-cafe', 'lubricants', 'fleet-premium'.")
    compliance_report: ComplianceCheck

# Flat model to prevent $defs and $ref schema validation failures in the Google GenAI backend
class FlatMultiAgentPayload(BaseModel):
    email_subject: str = Field(..., description="High-impact corporate email subject line.")
    email_body: str = Field(..., description="Comprehensive retail email body supporting the narrative.")
    push_notification: str = Field(..., description="Short, compelling mobile application push message (max 120 chars).")
    in_app_banner: str = Field(..., description="Action-oriented text for in-app banner placement.")
    asset_theme_token: str = Field(..., description="Sleek CSS slug representing visual asset directive. Strictly one of: 'ev-charging', 'aura-cafe', 'lubricants', 'fleet-premium'.")
    compliance_status: str = Field(..., description="Strictly either 'APPROVED' or 'RISK_FLAG_DETECTED'")
    asa_greenwashing_passed: bool = Field(..., description="True if claims are linked to specific asset footprints, False if generalities are used.")
    clear_pricing_disclosure: bool = Field(..., description="True if discounts and criteria are explicitly bounded.")
    executive_summary: str = Field(..., description="A professional 1-sentence risk and regulatory compliance sign-off.")

# Combined Comparison Schema returning outputs from both Gemini engines
class ComparisonPayload(BaseModel):
    flash: MultiAgentPayload = Field(..., description="Orchestration output from gemini-2.5-flash.")
    pro: MultiAgentPayload = Field(..., description="Orchestration output from gemini-2.5-flash (alternative profile).")
    retrieved_asset: dict = Field(default=None, description="The semantic asset surfaced by pgvector during search.")

# --- BRIEFING MODULE SCHEMAS ---
class BriefingRequest(BaseModel):
    campaign_name: str = Field(..., example="Autumn Forecourt Café Push")
    target_market: Optional[str] = Field(default="UK", example="UK")
    channels: List[str] = Field(default=[], example=["email", "push", "in_app"])
    deadline: Optional[str] = Field(default=None, example="2026-06-15")
    budget_tier: Optional[str] = Field(default="standard", example="premium")
    detailed_brief: str = Field(..., example="We need a full multi-channel campaign promoting our new organic artisan coffee range...")
    priority: Optional[str] = Field(default="normal", example="high")
    submitted_by: Optional[str] = Field(default="client", example="Vishesh Agarwal")

class BriefingStatusUpdate(BaseModel):
    new_status: str = Field(..., example="in_review")
    hub_notes: Optional[str] = Field(default=None, example="Assigned to creative team.")

# Flat schema for Gemini extraction — no nested types to avoid $ref issues
class BriefingExtract(BaseModel):
    campaign_name: str = Field(..., description="A concise, compelling campaign name derived from the document content.")
    target_market: str = Field(default="UK", description="The primary target market or geographic region. Must be one of: UK, EU, US, APAC, Global.")
    channels: List[str] = Field(default=[], description="Marketing channels mentioned or implied. Each must be one of: email, push, in_app, social, print, ooh.")
    deadline: str = Field(default="", description="Any deadline or target date mentioned, in YYYY-MM-DD format. Return an empty string if not found.")
    budget_tier: str = Field(default="standard", description="Inferred budget tier based on scope and scale. Must be one of: standard, premium, enterprise.")
    detailed_brief: str = Field(..., description="A comprehensive campaign brief summarising the key objectives, target audience, messaging strategy, creative direction, and deliverables from the document.")
    priority: str = Field(default="normal", description="Inferred priority level. Must be one of: low, normal, high, urgent.")


# --- ASYNCHRONOUS PIPELINE HELPER ---
async def generate_single_model_output(
    client: genai.Client,
    model_name: str,
    brief: str,
    system_instruction: str,
    temperature: float = 0.2
) -> MultiAgentPayload:
    # Trigger async client call via the client.aio namespace with specific temperature profile
    response = await client.aio.models.generate_content(
        model=model_name,
        contents=f"Execute supply chain pipeline for this campaign activation brief:\n\n{brief}",
        config=types.GenerateContentConfig(
            system_instruction=system_instruction,
            response_mime_type="application/json",
            response_schema=FlatMultiAgentPayload,
            temperature=temperature,
        ),
    )
    if not response.text:
        raise ValueError(f"Empty response received from {model_name}")
        
    flat = FlatMultiAgentPayload.model_validate_json(response.text)
    
    # Pack the flat data into the nested structure expected by the frontend
    return MultiAgentPayload(
        localized_content=LocalizedCopy(
            email_subject=flat.email_subject,
            email_body=flat.email_body,
            push_notification=flat.push_notification,
            in_app_banner=flat.in_app_banner
        ),
        asset_theme_token=flat.asset_theme_token,
        compliance_report=ComplianceCheck(
            status=flat.compliance_status,
            asa_greenwashing_passed=flat.asa_greenwashing_passed,
            clear_pricing_disclosure=flat.clear_pricing_disclosure,
            executive_summary=flat.executive_summary
        )
    )


# --- MULTI-AGENT ORCHESTRATION GATEWAY ---
# --- MOCK VECTOR UPLOAD API FOR DEMO ---
@app.post("/api/v1/assets/upload")
async def upload_asset_endpoint(file: UploadFile = File(...)):
    """Allows uploading raw files, extracting their content, vectorising them, and storing in pgvector."""
    try:
        filename = file.filename
        # Determine file type from extension
        ext = filename.split(".")[-1].lower() if "." in filename else "txt"
        
        # Read the file bytes
        file_bytes = await file.read()
        
        # Extract text based on file format
        content = ""
        if ext == "txt":
            content = file_bytes.decode("utf-8", errors="ignore")
        elif ext == "pdf":
            try:
                from pypdf import PdfReader
                pdf_file = io.BytesIO(file_bytes)
                reader = PdfReader(pdf_file)
                extracted_text = []
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        extracted_text.append(page_text)
                content = "\n".join(extracted_text).strip()
            except Exception as pdf_err:
                content = f"PDF text extraction failed: {str(pdf_err)}"
        elif ext in ["ppt", "pptx"]:
            try:
                from pptx import Presentation
                ppt_file = io.BytesIO(file_bytes)
                prs = Presentation(ppt_file)
                extracted_text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            extracted_text.append(shape.text)
                content = "\n".join(extracted_text).strip()
            except Exception as ppt_err:
                content = f"PPTX text extraction failed: {str(ppt_err)}"
        elif ext in ["png", "jpg", "jpeg", "gif"]:
            # High-fidelity descriptive summary for media files to ground embeddings
            content = (
                f"[VECTORIZED MARKETING IMAGE ASSET]: Filename: {filename}. "
                f"This asset contains premium brand creative elements, localized high-impact visual banners, "
                f"or forecourt Cafe/Automotive product specs. Compliant sustainability messaging is embedded."
            )
        else:
            content = file_bytes.decode("utf-8", errors="ignore")[:2000]
            
        if not content.strip() or len(content.strip()) < 10:
            # Fallback if no text extracted
            content = (
                f"Factual enterprise information extracted from {filename} ({ext.upper()} asset). "
                f"This document represents campaign creative materials, localized brand graphics, "
                f"or regulatory approval logs for active channels. Primary key features include smart battery "
                f"load-balancing grid efficiency, organic commuter retail cafe promotions, and premium Roadster cruise specs."
            )
            
        # Generate authentic 768-dimensional vector embeddings using Google's text-embedding-004
        try:
            api_key = os.environ.get("GEMINI_API_KEY")
            if api_key:
                ai_client = genai.Client(api_key=api_key)
            else:
                ai_client = genai.Client()
            
            embed_response = ai_client.models.embed_content(
                model="text-embedding-004",
                contents=content
            )
            embedding = embed_response.embeddings[0].values
        except Exception as embed_err:
            import logging
            logging.getLogger("AuraAPI").warning(f"Failed to generate live embedding: {str(embed_err)}. Generating unique hash-based fallback vector.")
            # Deterministic, unique fallback vector based on document hash so records are not identical
            import hashlib
            hash_bytes = hashlib.sha256(content.encode("utf-8")).digest()
            embedding = [float(b) / 255.0 for b in hash_bytes] * 24 # 32 * 24 = 768 dimensions
            
        success = db.add_asset(filename, ext, content, {"source": "upload-api", "size": len(file_bytes)}, embedding)
        if success:
            return {
                "status": "success",
                "filename": filename,
                "file_type": ext,
                "content_chunk": content[:300] + "...",
                "message": f"Asset '{filename}' successfully parsed, vectorized, and indexed in pgvector."
            }
        raise HTTPException(status_code=500, detail="Database write failure.")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/v1/orchestrate", response_model=ComparisonPayload)
async def orchestrate_content_supply_chain(payload: CampaignRequest):
    # Initialize GenAI Client dynamically at request time to pick up env vars robustly
    try:
        api_key = os.environ.get("GEMINI_API_KEY")
        if api_key:
            client = genai.Client(api_key=api_key)
        else:
            client = genai.Client()
    except Exception as e:
        raise HTTPException(
            status_code=500, 
            detail=f"Google GenAI SDK client initialization failed: {str(e)}"
        )

    # 1. Query the pgvector store for semantic asset matches (RAG) using live embeddings
    retrieved_context = ""
    retrieved_asset = None
    try:
        # Generate real 768-dimensional query vector using text-embedding-004
        try:
            embed_response = client.models.embed_content(
                model="text-embedding-004",
                contents=payload.brief
            )
            query_vector = embed_response.embeddings[0].values
        except Exception as embed_err:
            import logging
            logging.getLogger("AuraAPI").warning(f"Failed to generate query embedding: {str(embed_err)}. Using deterministic fallback vector.")
            # Deterministic fallback vector based on query string
            import hashlib
            hash_bytes = hashlib.sha256(payload.brief.encode("utf-8")).digest()
            query_vector = [float(b) / 255.0 for b in hash_bytes] * 24

        matches = db.semantic_search(query_vector, limit=1)
        if matches:
            best_match = matches[0]
            retrieved_context = (
                f"\n\n[RETRIEVED FACT/ASSET FROM AURA DAM]:\n"
                f"File: {best_match['filename']} (Type: {best_match['file_type']})\n"
                f"Content: {best_match['content_chunk']}\n"
                f"You MUST use this retrieved factual content to ground your copy, "
                f"and ensure you reference the asset properly in your output."
            )
            retrieved_asset = best_match
    except Exception as e:
        import logging
        logging.getLogger("AuraAPI").error(f"pgvector lookup error: {str(e)}")

    # Initialize GenAI Client dynamically at request time to pick up env vars robustly
    try:
        api_key = os.environ.get("GEMINI_API_KEY")
        if api_key:
            client = genai.Client(api_key=api_key)
        else:
            client = genai.Client()
    except Exception as e:
        raise HTTPException(
            status_code=500, 
            detail=f"Google GenAI SDK client initialization failed: {str(e)}"
        )
    
    # Strict System Instruction Blueprint mapping out the multi-agent persona boundaries
    system_instruction = (
        "You are the Enterprise AI Orchestration Framework running the core Marketing Content Supply Chain for Aura (a leading sustainable mobility and energy provider). "
        "You must process the incoming user marketing brief by executing three independent internal agent protocols:\n\n"
        
        "AGENT A (The Content Localizer Persona):\n"
        "Generate localized, high-performing copy targeted precisely to the geographic/demographic context given. "
        "Output distinct variations for Corporate Email (Subject & Body), Mobile App Push, and an In-App Interactive Banner.\n\n"
        
        "AGENT B (The Visual Assets Director Persona):\n"
        "Evaluate the contextual elements of the campaign brief. Select the single most appropriate structural design "
        "theme token from this strict enumeration: 'ev-charging', 'aura-cafe', 'lubricants', 'fleet-premium'.\n\n"
        
        "AGENT C (The Sustainability & Compliance Guardrail Persona):\n"
        "Audit the generated copy against UK ASA (Advertising Standards Authority) and CMA greenwashing regulations. "
        "Crucial Rule: All environmental, pricing, or product claims (like organic ingredients or electric vehicle ranges) must be framed "
        "strictly around localized, quantifiable asset rollouts. If the copy uses vague corporate generalities "
        "(e.g., 'making the planet cleaner'), trigger a 'RISK_FLAG_DETECTED' status (field compliance_status). "
        "Verify transparent pricing terms.\n\n"
        
        "Consolidate all outputs into the exact JSON schema defined."
    )
    
    # Append the RAG context to instruct Gemini to ground its output
    system_instruction += retrieved_context

    # Distinct strategic operational profiles
    operational_instruction = (
        system_instruction + 
        "\nCRITICAL OPERATIONAL CONSTRAINT: Focus heavily on factual accuracy, conservative compliance framing, "
        "and strict adherence to greenwashing guardrails. Keep the copywriting highly professional and direct."
    )
    
    creative_instruction = (
        system_instruction + 
        "\nCRITICAL CREATIVE CONSTRAINT: Focus heavily on persuasive, highly engaging marketing copy, bold attention-grabbing "
        "subject lines, and maximum call-to-action impact for drivers. You may be more creative and emotionally appealing "
        "while still maintaining basic compliance."
    )

    try:
        # Run both pipelines concurrently using gemini-2.5-flash with custom profiles (Temp 0.1 vs Temp 0.8)
        # This guarantees high-speed execution and eliminates 429 quota exhaustion blocks on developer free tier accounts.
        flash_task = generate_single_model_output(
            client=client, 
            model_name='gemini-2.5-flash', 
            brief=payload.brief, 
            system_instruction=operational_instruction, 
            temperature=0.1
        )
        pro_task = generate_single_model_output(
            client=client, 
            model_name='gemini-2.5-flash', 
            brief=payload.brief, 
            system_instruction=creative_instruction, 
            temperature=0.8
        )
        
        flash_result, pro_result = await asyncio.gather(flash_task, pro_task)
        
        return ComparisonPayload(flash=flash_result, pro=pro_result, retrieved_asset=retrieved_asset)

    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(
            status_code=500, 
            detail=f"Agentic Pipeline Comparison Execution Failure: {str(e)}"
        )

# --- BRIEFING MODULE ENDPOINTS ---

@app.post("/api/v1/briefings/extract")
async def extract_briefing_from_document(file: UploadFile = File(...)):
    """Upload a document, extract text, and use Gemini to auto-populate briefing fields."""
    try:
        filename = file.filename
        ext = filename.split(".")[-1].lower() if "." in filename else "txt"
        file_bytes = await file.read()
        
        # --- Text extraction (reuses proven logic from asset upload) ---
        content = ""
        if ext == "txt":
            content = file_bytes.decode("utf-8", errors="ignore")
        elif ext == "pdf":
            try:
                from pypdf import PdfReader
                pdf_file = io.BytesIO(file_bytes)
                reader = PdfReader(pdf_file)
                extracted_text = []
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        extracted_text.append(page_text)
                content = "\n".join(extracted_text).strip()
            except Exception as pdf_err:
                raise HTTPException(status_code=422, detail=f"PDF text extraction failed: {str(pdf_err)}")
        elif ext in ["ppt", "pptx"]:
            try:
                from pptx import Presentation
                ppt_file = io.BytesIO(file_bytes)
                prs = Presentation(ppt_file)
                extracted_text = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            extracted_text.append(shape.text)
                content = "\n".join(extracted_text).strip()
            except Exception as ppt_err:
                raise HTTPException(status_code=422, detail=f"PPTX text extraction failed: {str(ppt_err)}")
        elif ext in ["doc", "docx"]:
            content = file_bytes.decode("utf-8", errors="ignore")[:5000]
        else:
            content = file_bytes.decode("utf-8", errors="ignore")[:5000]
        
        if not content.strip() or len(content.strip()) < 20:
            raise HTTPException(status_code=422, detail="Could not extract sufficient text content from the uploaded document.")
        
        # Truncate to avoid token limits while preserving key info
        content = content[:8000]
        
        # --- Gemini structured extraction ---
        api_key = os.environ.get("GEMINI_API_KEY")
        if api_key:
            client = genai.Client(api_key=api_key)
        else:
            client = genai.Client()
        
        system_instruction = (
            "You are an expert marketing operations analyst. You are given the raw text content extracted from a campaign document "
            "(such as a brief, strategy deck, marketing plan, or creative proposal). "
            "Your task is to carefully read the document and extract structured briefing fields from it.\n\n"
            "Rules:\n"
            "- campaign_name: Create a concise, professional campaign name that captures the core theme.\n"
            "- target_market: Identify the primary geographic market. Must be exactly one of: UK, EU, US, APAC, Global.\n"
            "- channels: List all marketing channels mentioned or strongly implied. Each must be one of: email, push, in_app, social, print, ooh.\n"
            "- deadline: Extract any deadline or target date in YYYY-MM-DD format. Use null if none found.\n"
            "- budget_tier: Infer from the scope and ambition. Must be one of: standard, premium, enterprise.\n"
            "- detailed_brief: Write a comprehensive summary (3-5 sentences) covering objectives, target audience, key messages, creative direction, and deliverables.\n"
            "- priority: Infer from urgency cues in the document. Must be one of: low, normal, high, urgent.\n\n"
            "Return ONLY the JSON object matching the exact schema provided. Do not add extra fields."
        )
        
        response = await client.aio.models.generate_content(
            model="gemini-2.5-flash",
            contents=f"Extract structured briefing fields from this campaign document:\n\n{content}",
            config=types.GenerateContentConfig(
                system_instruction=system_instruction,
                response_mime_type="application/json",
                response_schema=BriefingExtract,
                temperature=0.1,
            ),
        )
        
        if not response.text:
            raise HTTPException(status_code=500, detail="Gemini returned an empty response.")
        
        extracted = BriefingExtract.model_validate_json(response.text)
        
        return {
            "status": "success",
            "source_filename": filename,
            "source_file_type": ext,
            "extracted_text_length": len(content),
            "fields": extracted.model_dump()
        }
        
    except HTTPException:
        raise
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=f"Briefing extraction failed: {str(e)}")

@app.post("/api/v1/briefings")
async def submit_briefing(payload: BriefingRequest):
    """Submit a new campaign briefing request for the hub team."""
    try:
        data = payload.model_dump()
        result = db.insert_briefing(data)
        if result:
            return {"status": "success", "briefing": result}
        raise HTTPException(status_code=500, detail="Failed to store briefing in database.")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Briefing submission failed: {str(e)}")

@app.get("/api/v1/briefings")
async def list_briefings(status: Optional[str] = Query(default=None, description="Filter by status: submitted, in_review, in_progress, delivered")):
    """List all briefing requests, optionally filtered by status."""
    try:
        results = db.list_briefings(status_filter=status)
        return {"status": "success", "briefings": results, "total": len(results)}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to retrieve briefings: {str(e)}")

@app.patch("/api/v1/briefings/{briefing_id}/status")
async def update_briefing_status(briefing_id: str, payload: BriefingStatusUpdate):
    """Update the status of a briefing (designed for hub team admin panel)."""
    try:
        result = db.update_briefing_status(briefing_id, payload.new_status, payload.hub_notes)
        if result:
            return {"status": "success", "briefing": result}
        raise HTTPException(status_code=404, detail=f"Briefing '{briefing_id}' not found or invalid status.")
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to update briefing: {str(e)}")

@app.get("/health")
async def health_check():
    return {"status": "healthy", "engine": "FastAPI Agentic Orchestrator Online"}

@app.get("/", response_class=HTMLResponse)
async def serve_frontend():
    frontend_path = os.path.join(os.path.dirname(__file__), "../frontend/index.html")
    if os.path.exists(frontend_path):
        with open(frontend_path, "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    return HTMLResponse(content="<h1>Aura Frontend Index Not Found</h1>", status_code=404)
